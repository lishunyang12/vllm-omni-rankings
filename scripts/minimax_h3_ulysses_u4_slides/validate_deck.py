#!/usr/bin/env python3
"""Structural and conservative text-capacity checks for the generated deck."""

import math
import sys
from pathlib import Path

from pptx import Presentation


def text_weight(text: str) -> float:
    return sum(1.0 if ord(ch) > 127 else 0.58 for ch in text)


def main(path: str) -> int:
    deck = Path(path)
    prs = Presentation(deck)
    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400
    warnings = []
    text_shapes = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape_idx, shape in enumerate(slide.shapes, 1):
            x = shape.left / 914400
            y = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            if x < -0.001 or y < -0.001 or x + w > slide_w + 0.001 or y + h > slide_h + 0.001:
                warnings.append(f"slide {slide_idx} shape {shape_idx}: out of bounds")

            if not shape.has_text_frame or not shape.text.strip():
                continue
            text_shapes += 1
            tf = shape.text_frame
            usable_w_pt = max(1.0, (w - 0.02) * 72)
            usable_h_pt = max(1.0, (h - 0.02) * 72)
            needed_h_pt = 0.0
            for paragraph in tf.paragraphs:
                font_sizes = [run.font.size.pt for run in paragraph.runs if run.font.size]
                font_size = max(font_sizes) if font_sizes else 18.0
                chars_per_line = max(1.0, usable_w_pt / font_size)
                lines = max(1, math.ceil(text_weight(paragraph.text) / chars_per_line))
                needed_h_pt += lines * font_size * 1.08
            if needed_h_pt > usable_h_pt * 1.12:
                sample = shape.text.replace("\n", " / ")[:72]
                warnings.append(
                    f"slide {slide_idx} shape {shape_idx}: possible text overflow "
                    f"({needed_h_pt:.0f}pt > {usable_h_pt:.0f}pt): {sample}"
                )

    print(f"deck={deck}")
    print(f"slides={len(prs.slides)} size={slide_w:.3f}x{slide_h:.3f}in text_shapes={text_shapes}")
    if warnings:
        print(f"warnings={len(warnings)}")
        for warning in warnings:
            print(warning)
        return 1
    print("warnings=0")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1]))
