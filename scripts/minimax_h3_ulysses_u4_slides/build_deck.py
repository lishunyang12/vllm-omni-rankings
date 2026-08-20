#!/usr/bin/env python3
"""Generate the MiniMax-H3 Ulysses academic slide deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "MiniMax-H3_Ulysses_U4_学术解析.pptx"

SW, SH = 13.333, 7.5
BG = "07141F"
PANEL = "102633"
PANEL_2 = "16313F"
INK = "F4F8FA"
MUTED = "A8BBC7"
TEAL = "2DD4BF"
BLUE = "4EA8DE"
ORANGE = "F59E0B"
PURPLE = "A78BFA"
RED = "FB7185"
GREEN = "86EFAC"
GRID = "274656"
FONT = "PingFang SC"
MONO = "Menlo"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)
    return slide


def add_rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=20,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.01,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraphs = str(text).split("\n")
    for idx, value in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = value
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, kicker, title, subtitle=None):
    add_text(slide, 0.62, 0.34, 12.0, 0.28, kicker.upper(), 10, TEAL, True, MONO)
    add_text(slide, 0.62, 0.68, 12.0, 0.55, title, 27, INK, True)
    add_rect(slide, 0.62, 1.34, 1.15, 0.035, TEAL, TEAL, radius=False)
    if subtitle:
        add_text(slide, 1.95, 1.23, 10.7, 0.30, subtitle, 11, MUTED)


def add_footer(slide, number, source="本地源码与实验记录；分析日期 2026-08-20"):
    add_rect(slide, 0.62, 7.08, 12.05, 0.012, GRID, GRID, radius=False)
    add_text(slide, 0.62, 7.13, 11.2, 0.20, source, 8, MUTED, font=MONO)
    add_text(slide, 12.0, 7.11, 0.66, 0.22, f"{number:02d}", 9, TEAL, True, MONO, PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, items, size=17, color=INK, bullet_color=TEAL, gap=0.52):
    for i, item in enumerate(items):
        yy = y + i * gap
        add_rect(slide, x, yy + 0.14, 0.08, 0.08, bullet_color, bullet_color, radius=False)
        add_text(slide, x + 0.20, yy, w - 0.20, gap, item, size, color)


def add_pill(slide, x, y, w, text, fill=TEAL, color=BG, size=11):
    add_rect(slide, x, y, w, 0.34, fill, fill, radius=True)
    add_text(slide, x, y + 0.01, w, 0.30, text, size, color, True, MONO, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def add_metric(slide, x, y, w, value, label, accent=TEAL, detail=None):
    add_rect(slide, x, y, w, 1.16, PANEL, GRID)
    add_text(slide, x + 0.18, y + 0.12, w - 0.36, 0.48, value, 25, accent, True, MONO)
    add_text(slide, x + 0.18, y + 0.63, w - 0.36, 0.28, label, 12, INK, True)
    if detail:
        add_text(slide, x + 0.18, y + 0.91, w - 0.36, 0.20, detail, 8, MUTED)


def add_code(slide, x, y, w, h, code, size=12, accent=TEAL):
    add_rect(slide, x, y, w, h, "0B1D28", GRID)
    add_rect(slide, x, y, 0.055, h, accent, accent, radius=False)
    add_text(slide, x + 0.18, y + 0.15, w - 0.34, h - 0.25, code, size, "D8E5EC", font=MONO, line_spacing=0.92)


def add_chevron(slide, x, y, w=0.38, h=0.48, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    return shape


def add_process_box(slide, x, y, w, h, title, body, accent=TEAL, title_size=15, body_size=11):
    add_rect(slide, x, y, w, h, PANEL, GRID)
    add_rect(slide, x, y, 0.06, h, accent, accent, radius=False)
    add_text(slide, x + 0.18, y + 0.13, w - 0.30, 0.32, title, title_size, accent, True)
    add_text(slide, x + 0.18, y + 0.52, w - 0.30, h - 0.62, body, body_size, INK)


def add_rank(slide, x, y, rank, seq, heads, color):
    add_rect(slide, x, y, 2.58, 0.72, PANEL, color)
    add_pill(slide, x + 0.12, y + 0.18, 0.62, f"R{rank}", color, BG, 10)
    add_text(slide, x + 0.87, y + 0.10, 1.56, 0.24, seq, 10, INK, True, MONO)
    add_text(slide, x + 0.87, y + 0.39, 1.56, 0.20, heads, 9, MUTED, font=MONO)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    prs.core_properties.title = "MiniMax-H3 负载中的 Ulysses Sequence Parallelism"
    prs.core_properties.subject = "H3 U4 execution path, communication model, profiling and optimization"
    prs.core_properties.author = "Codex × Local source analysis"

    # 01 — title
    slide = blank_slide(prs)
    add_pill(slide, 0.70, 0.60, 2.20, "SYSTEMS · ACADEMIC", TEAL)
    add_text(slide, 0.70, 1.45, 11.7, 1.25, "MiniMax-H3 负载中的\nUlysses Sequence Parallelism", 35, INK, True)
    add_text(slide, 0.74, 3.05, 10.8, 0.62, "从 packed multimodal sequence 到 Q/K/V All-to-All、反向重分片与最终 AllGather", 18, MUTED)
    add_rect(slide, 0.74, 4.25, 11.80, 1.55, PANEL, GRID)
    add_text(slide, 1.05, 4.57, 3.1, 0.30, "研究对象", 11, TEAL, True, MONO)
    add_text(slide, 1.05, 4.96, 3.1, 0.52, "H3 DiT · U4 · TP1", 19, INK, True)
    add_text(slide, 4.72, 4.57, 3.1, 0.30, "核心尺度", 11, ORANGE, True, MONO)
    add_text(slide, 4.72, 4.96, 3.1, 0.52, "S=37,760 · H=5,376", 19, INK, True, MONO)
    add_text(slide, 8.65, 4.57, 3.1, 0.30, "证据基础", 11, PURPLE, True, MONO)
    add_text(slide, 8.65, 4.96, 3.1, 0.52, "源码 + B300 实验", 21, INK, True)
    add_text(slide, 0.74, 6.50, 11.8, 0.28, "agent/minimax-h3-online-fp8 @ f76f8e58fb  ·  2026-08-20", 10, MUTED, font=MONO)
    add_footer(slide, 1)

    # 02 — questions and conclusions
    slide = blank_slide(prs)
    add_title(slide, "01 · Research framing", "研究问题与核心结论", "先给出系统层答案，再逐层证明")
    add_rect(slide, 0.62, 1.70, 5.70, 4.95, PANEL, GRID)
    add_text(slide, 0.92, 1.98, 4.9, 0.38, "三个研究问题", 18, TEAL, True)
    add_bullets(slide, 0.95, 2.62, 4.95, [
        "H3 在哪里进入、退出 sequence-parallel 区域？",
        "为什么每个 DiT block 有 3×QKV + 1×Output All-to-All？",
        "为什么末端仍有 AllGather；如何把 payload 缩小 42×？",
    ], 16, gap=0.98)
    add_rect(slide, 6.58, 1.70, 6.10, 4.95, PANEL_2, GRID)
    add_text(slide, 6.90, 1.98, 5.3, 0.38, "四条结论", 18, ORANGE, True)
    add_bullets(slide, 6.92, 2.52, 5.18, [
        "_sp_plan 只定义边界；Ulysses 真正执行在通用 Attention strategy。",
        "核心是 S-shard ↔ H-shard 的分布式转置，而非简单 AllGather K/V。",
        "U4 每次 transformer forward 有 200 次 All-to-All；5 sigma 点对应 4 次 forward。",
        "当前分支先 gather [S,5376]；严格边界优化应先本地投影到 [S/4,128]。",
    ], 15, ORANGE, gap=0.84)
    add_footer(slide, 2)

    # 03 — baseline
    slide = blank_slide(prs)
    add_title(slide, "02 · Experimental baseline", "可复现实验基线", "所有计数与尺寸均对应同一 U4 配置")
    add_metric(slide, 0.62, 1.72, 2.86, "4× B300", "GPU", TEAL, "267.68 GiB / GPU")
    add_metric(slide, 3.65, 1.72, 2.86, "TP1 · U4", "DiT parallelism", BLUE, "Ring degree = 1")
    add_metric(slide, 6.68, 1.72, 2.86, "768×1344", "输出分辨率", ORANGE, "124 frames @ 24 fps")
    add_metric(slide, 9.71, 1.72, 2.96, "5 σ points", "采样设置", PURPLE, "4 denoise transitions")
    add_rect(slide, 0.62, 3.18, 7.26, 3.38, PANEL, GRID)
    add_text(slide, 0.92, 3.45, 6.6, 0.32, "运行配置", 17, TEAL, True)
    add_code(slide, 0.92, 3.94, 6.64, 1.54,
             "tp1_ulysses4_ring1\ntext_encoder_tp4_vae_tile4\nattention_backend = CUDNN_ATTN\nprecision = checkpoint BF16/FP32", 14)
    add_text(slide, 0.94, 5.70, 6.50, 0.56, "注意：Text Encoder TP4 与 VAE tile4 是外围并行，本文的 Ulysses 分析聚焦 DiT denoise hot path。", 12, MUTED)
    add_rect(slide, 8.15, 3.18, 4.53, 3.38, PANEL_2, GRID)
    add_text(slide, 8.46, 3.45, 3.9, 0.32, "观测到的阶段耗时", 17, ORANGE, True)
    add_text(slide, 8.46, 4.05, 1.55, 0.26, "T2VA", 12, MUTED, True, MONO)
    add_text(slide, 9.94, 3.95, 2.0, 0.40, "4.559 s", 22, INK, True, MONO)
    add_text(slide, 8.46, 4.62, 1.55, 0.26, "FL2VA", 12, MUTED, True, MONO)
    add_text(slide, 9.94, 4.52, 2.0, 0.40, "3.310 s", 22, INK, True, MONO)
    add_text(slide, 8.46, 5.20, 3.65, 0.76, "以上为 diffuse 阶段墙钟时间；本 deck 不将其误解为纯通信时间。通信结论来自源码归因与尺寸模型。", 11, MUTED)
    add_footer(slide, 3, "results/h3-b300-tp1-u4-steps5/{summary.json,run.log}")

    # 04 — architecture
    slide = blank_slide(prs)
    add_title(slide, "03 · Workload anatomy", "H3 DiT 的结构尺度", "Ulysses 的收益与代价由 S、head 数和 block 数共同决定")
    metrics = [
        ("50", "DiT blocks", TEAL),
        ("5,376", "hidden size", BLUE),
        ("56 × 128", "heads × head dim", ORANGE),
        ("14,336", "FFN hidden", PURPLE),
    ]
    for i, (v, label, c) in enumerate(metrics):
        add_metric(slide, 0.62 + i * 3.03, 1.72, 2.86, v, label, c)
    add_process_box(slide, 0.62, 3.20, 2.55, 2.68, "Token refiner", "2 layers\n完整 sequence\nskip_sequence_parallel=True", PURPLE, 16, 13)
    add_chevron(slide, 3.30, 4.22)
    add_process_box(slide, 3.85, 3.20, 5.12, 2.68, "Main DiT stack", "50 × {AdaLN → Self-Attention → MLP}\nUlysses 仅在这个 sharded region 内生效", TEAL, 17, 14)
    add_chevron(slide, 9.12, 4.22)
    add_process_box(slide, 9.67, 3.20, 3.01, 2.68, "Final layer", "video width = 96\naudio width = 32\ncombined = 128", ORANGE, 16, 13)
    add_text(slide, 0.70, 6.19, 11.7, 0.42, "关键观察：当前分支在 5,376→128 的投影之前退出 SP，这是末端大 AllGather 的根因。", 15, INK, True)
    add_footer(slide, 4, "bench_inputs/MiniMax-H3/FL2VA/transformer/config.json")

    # 05 — packed sequence
    slide = blank_slide(prs)
    add_title(slide, "04 · Packed sequence", "多模态 token 如何打包成一个序列", "H3 不是 joint-attention 分支，而是把所有模态放进同一 packed sequence")
    x0, y0, total_w = 0.72, 2.05, 11.88
    segments = [
        ("text", "25", 1.45, PURPLE),
        ("img/vid cond", "0 (T2VA)", 1.75, BLUE),
        ("audio", "414", 1.65, ORANGE),
        ("video target", "37,296", 5.55, TEAL),
        ("pad", "25", 1.48, MUTED),
    ]
    x = x0
    for name, count, width, color in segments:
        add_rect(slide, x, y0, width, 1.15, color, BG, radius=False)
        add_text(slide, x + 0.05, y0 + 0.22, width - 0.10, 0.32, name, 13, BG, True, MONO, PP_ALIGN.CENTER)
        add_text(slide, x + 0.05, y0 + 0.61, width - 0.10, 0.26, count, 12, BG, True, MONO, PP_ALIGN.CENTER)
        x += width
    add_text(slide, 0.72, 3.54, 11.8, 0.48, "layout = [ text | imgvid_cond | audio | video_target | alignment pad ]", 18, INK, True, MONO, PP_ALIGN.CENTER)
    add_rect(slide, 0.72, 4.36, 11.88, 1.62, PANEL, GRID)
    add_text(slide, 1.02, 4.67, 3.15, 0.34, "为什么需要 pad？", 16, ORANGE, True)
    add_text(slide, 1.02, 5.10, 3.20, 0.47, "将总长度对齐到 64，保证 U4 等分和后端友好的 tile。", 13, MUTED)
    add_text(slide, 4.62, 4.67, 3.10, 0.34, "packed attention metadata", 15, TEAL, True, MONO)
    add_text(slide, 4.62, 5.10, 3.15, 0.47, "cu_seqlens = [0, 37,735, 37,760]\npad 作为第二个 document。", 12, INK, font=MONO)
    add_text(slide, 8.25, 4.67, 3.40, 0.34, "Ulysses 后的语义", 16, PURPLE, True)
    add_text(slide, 8.25, 5.06, 3.40, 0.62, "每个 head-owner rank 得到完整 packed sequence，metadata 仍使用全局边界。", 12, MUTED)
    add_footer(slide, 5, "models/minimax_h3/packed_sequence.py")

    # 06 — length derivation
    slide = blank_slide(prs)
    add_title(slide, "05 · Shape derivation", "T2VA 的 S=37,760 如何得到", "显式推导，避免把逻辑序列长度与 padded 长度混淆")
    add_process_box(slide, 0.62, 1.72, 3.42, 1.42, "Video latent", "T=124 → latent_t=37\nH=768→48，W=1344→84", TEAL, 16, 13)
    add_chevron(slide, 4.18, 2.20)
    add_process_box(slide, 4.72, 1.72, 3.42, 1.42, "Spatial patch", "patch=2×2\n每帧 24×42 = 1,008 rows", BLUE, 16, 13)
    add_chevron(slide, 8.29, 2.20)
    add_process_box(slide, 8.83, 1.72, 3.84, 1.42, "Video rows", "37 × 1,008\n= 37,296", ORANGE, 16, 15)
    add_rect(slide, 0.62, 3.55, 12.05, 2.42, PANEL, GRID)
    add_text(slide, 0.95, 3.88, 11.4, 0.50, "used = text + audio + video = 25 + 414 + 37,296 = 37,735", 21, INK, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 0.95, 4.55, 11.4, 0.50, "S = align₆₄(37,735) = 37,760  →  padding = 25", 22, TEAL, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 0.95, 5.25, 11.4, 0.36, "U4 local sequence = S / 4 = 9,440 rows per rank", 18, ORANGE, True, MONO, PP_ALIGN.CENTER)
    add_footer(slide, 6, "packed_sequence.py + time_request.py；T2VA text presentation length=25")

    # 07 — topology
    slide = blank_slide(prs)
    add_title(slide, "06 · Parallel topology", "TP1 · Ulysses4 · Ring1 的进程组", "四张卡属于同一个 Ulysses group；每个 Ring group 是单例")
    colors = [TEAL, BLUE, ORANGE, PURPLE]
    for i in range(4):
        add_rank(slide, 0.76 + i * 3.08, 1.88, i, f"seq = S/4", "heads = 56", colors[i])
    add_rect(slide, 0.76, 2.92, 11.82, 0.64, PANEL_2, TEAL)
    add_text(slide, 0.90, 3.08, 11.50, 0.26, "ulysses_group = [0, 1, 2, 3]   ·   sp_size = 4", 16, TEAL, True, MONO, PP_ALIGN.CENTER)
    for i in range(4):
        add_rect(slide, 0.76 + i * 3.08, 3.88, 2.58, 0.58, PANEL, GRID)
        add_text(slide, 0.86 + i * 3.08, 4.03, 2.38, 0.22, f"ring_group = [{i}]", 11, MUTED, True, MONO, PP_ALIGN.CENTER)
    add_rect(slide, 0.76, 4.86, 11.82, 1.22, PANEL, GRID)
    add_text(slide, 1.02, 5.12, 3.16, 0.28, "序列并行前", 14, ORANGE, True)
    add_text(slide, 1.02, 5.42, 3.16, 0.45, "各 rank：S/4 tokens × 56 heads", 12, INK, font=MONO)
    add_text(slide, 4.85, 5.12, 3.16, 0.28, "Attention 内部", 14, TEAL, True)
    add_text(slide, 4.85, 5.47, 3.16, 0.30, "各 rank：S tokens × 14 heads", 13, INK, font=MONO)
    add_text(slide, 8.62, 5.12, 3.16, 0.28, "本配置没有", 14, PURPLE, True)
    add_text(slide, 8.62, 5.47, 3.16, 0.30, "Ring P2P / DiT TP collectives", 13, INK, font=MONO)
    add_footer(slide, 7, "run.log:45–52；parallel_state.py:620–682")

    # 08 — sp plan
    slide = blank_slide(prs)
    add_title(slide, "07 · SP boundary", "_sp_plan：用 hook 声明进入与退出边界", "0/1/2 是 tuple 位置；split_dim=0 才是 tensor 维度")
    code = """_sp_plan = {
  \"sp_prepare\": {
    0: SequenceParallelInput(split_dim=0,
         expected_dims=2, split_output=True),
    1: SequenceParallelInput(split_dim=0,
         expected_dims=2, split_output=True),
    2: SequenceParallelInput(split_dim=0,
         expected_dims=1, split_output=True),
  },
  \"sp_gather\": SequenceParallelOutput(
       gather_dim=0, expected_dims=2),
}"""
    add_code(slide, 0.62, 1.70, 6.22, 4.94, code, 11)
    add_process_box(slide, 7.16, 1.70, 5.51, 1.28, "sp_prepare", "passthrough module 返回三元 tuple；post-hook 对三个输出沿 S 同步切片。", TEAL, 16, 13)
    add_process_box(slide, 7.16, 3.18, 5.51, 1.28, "SP state", "实际发生切分后：_sp_shard_depth += 1，因此 Attention strategy 激活。", ORANGE, 16, 13)
    add_process_box(slide, 7.16, 4.66, 5.51, 1.28, "sp_gather", "对 hidden 输出执行 AllGather(dim=0)，随后 _sp_shard_depth -= 1。", PURPLE, 16, 13)
    add_text(slide, 7.20, 6.12, 5.35, 0.48, "边界定义 ≠ Ulysses 算法；真正 A2A 在 Attention strategy。", 12, INK, True)
    add_footer(slide, 8, "minimax_h3_transformer.py:743–759,785–804；sequence_parallel.py")

    # 09 — prepare lifecycle
    slide = blank_slide(prs)
    add_title(slide, "08 · Entering SP", "sp_prepare 调用前后发生了什么", "本地切片不产生 NCCL 通信，但它打开后续 Ulysses 的动态作用域")
    add_process_box(slide, 0.62, 1.86, 3.34, 3.52, "调用前：replicated", "hidden [37,760, 5,376]\nrope   [37,760, …]\ncombined [37,760]\n\n每个 rank 持有相同完整输入", PURPLE, 17, 14)
    add_chevron(slide, 4.18, 3.24, 0.50, 0.72, ORANGE)
    add_process_box(slide, 4.91, 1.86, 3.46, 3.52, "sp_prepare hook", "三个 tensor 使用相同 rank slice\n\nrank r rows:\n[r·9,440 : (r+1)·9,440]\n\nctx.sp_active = True", ORANGE, 17, 14)
    add_chevron(slide, 8.61, 3.24, 0.50, 0.72, TEAL)
    add_process_box(slide, 9.35, 1.86, 3.32, 3.52, "调用后：sharded", "hidden [9,440, 5,376]\nrope   [9,440, …]\ncombined [9,440]\n\n进入 50 层 local block stack", TEAL, 17, 14)
    add_rect(slide, 0.62, 5.76, 12.05, 0.67, PANEL, GRID)
    add_text(slide, 0.88, 5.94, 11.55, 0.28, "cu_seqlens / max_seqlen 保持全局语义：A2A 后每个 head-owner rank 会看到完整 S。", 15, INK, True, FONT, PP_ALIGN.CENTER)
    add_footer(slide, 9, "minimax_h3_transformer.py:1098–1123；hooks/sequence_parallel.py:236–268")

    # 10 — call chain
    slide = blank_slide(prs)
    add_title(slide, "09 · Call-path attribution", "从 H3 forward 到 NCCL kernel 的完整调用链", "代码归因比仅凭 profiler kernel 名称更可靠")
    chain = [
        ("H3 model", "sp_prepare → for block in blocks", TEAL),
        ("H3 block", "self.attn(h, rope, metadata)", BLUE),
        ("H3 attention", "QKV proj → norm/RoPE → Attention", ORANGE),
        ("generic layer", "strategy.pre_attention / post_attention", PURPLE),
        ("Ulysses", "SeqAllToAll4D.apply", RED),
        ("distributed comm", "dist.all_to_all_single", GREEN),
    ]
    y = 1.68
    for i, (name, body, color) in enumerate(chain):
        add_rect(slide, 0.92 + i * 1.84, y + i * 0.70, 3.04, 0.76, PANEL, color)
        add_text(slide, 1.08 + i * 1.84, y + 0.12 + i * 0.70, 1.15, 0.22, name, 8, color, True, MONO)
        add_text(slide, 2.25 + i * 1.84, y + 0.11 + i * 0.70, 1.50, 0.42, body, 9, INK, font=MONO)
    add_rect(slide, 0.62, 6.05, 12.05, 0.48, PANEL_2, GRID)
    add_text(slide, 0.88, 6.17, 11.55, 0.20, "关键断点：ulysses.py:332–334（Q/K/V）· ulysses.py:473（Output）· comm.py:51/86（collective）", 11, TEAL, True, MONO, PP_ALIGN.CENTER)
    add_footer(slide, 10, "minimax_h3_transformer.py → attention/layer.py → attention/parallel/ulysses.py → distributed/comm.py")

    # 11 — block flow
    slide = blank_slide(prs)
    add_title(slide, "10 · One DiT block", "单个 H3 block 的计算—通信时序", "MLP 保持 sequence-sharded；只有 Attention 内部临时切换为 head-sharded")
    nodes = [
        ("AdaLN + Norm", "[S/4,5376]", PURPLE),
        ("QKV proj", "3×[S/4,56,128]", BLUE),
        ("Q/K/V A2A", "3 collectives", ORANGE),
        ("Attention", "[S,14,128]", TEAL),
        ("Output A2A", "1 collective", ORANGE),
        ("OutProj + MLP", "[S/4,5376]", PURPLE),
    ]
    for i, (title, body, color) in enumerate(nodes):
        x = 0.44 + i * 2.12
        add_process_box(slide, x, 2.32, 1.75, 1.72, title, body, color, 12, 10)
        if i < len(nodes) - 1:
            add_chevron(slide, x + 1.80, 2.95, 0.26, 0.44, color)
    add_rect(slide, 0.62, 4.58, 12.05, 1.34, PANEL, GRID)
    add_text(slide, 0.90, 4.86, 3.20, 0.30, "通信区域", 14, ORANGE, True)
    add_text(slide, 0.90, 5.25, 3.20, 0.28, "4 × All-to-All / block", 16, INK, True, MONO)
    add_text(slide, 4.56, 4.86, 3.20, 0.30, "计算区域", 14, TEAL, True)
    add_text(slide, 4.56, 5.20, 3.20, 0.44, "full S × 14 heads / rank", 14, INK, True, MONO)
    add_text(slide, 8.46, 4.86, 3.20, 0.30, "输出布局", 14, PURPLE, True)
    add_text(slide, 8.46, 5.25, 3.20, 0.28, "恢复 S/4，下一层无需 gather", 15, INK, True)
    add_footer(slide, 11, "minimax_h3_transformer.py:638–681；attention/layer.py:263–290")

    # 12 — distributed transpose
    slide = blank_slide(prs)
    add_title(slide, "11 · Ulysses mathematics", "核心操作：分布式的 S/H 转置", "保持每 rank 元素数不变，改变谁拥有 token 与 head")
    add_rect(slide, 0.62, 1.80, 5.08, 3.95, PANEL, GRID)
    add_pill(slide, 0.92, 2.10, 1.30, "BEFORE", PURPLE)
    add_text(slide, 0.92, 2.68, 4.40, 0.52, "Xᵣ ∈ ℝᴮ×ˢ⁄ᴾ×ᴴ×ᴰ", 22, INK, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 1.12, 3.43, 4.0, 1.35, "rank r 拥有：\n• 一段 sequence\n• 全部 attention heads", 17, MUTED)
    add_text(slide, 0.92, 5.05, 4.40, 0.34, "[1, 9,440, 56, 128]", 17, PURPLE, True, MONO, PP_ALIGN.CENTER)
    add_chevron(slide, 5.96, 3.35, 1.22, 0.82, ORANGE)
    add_text(slide, 5.80, 4.30, 1.55, 0.62, "All-to-All\nscatter H · gather S", 10, ORANGE, True, MONO, PP_ALIGN.CENTER)
    add_rect(slide, 7.52, 1.80, 5.15, 3.95, PANEL_2, GRID)
    add_pill(slide, 7.82, 2.10, 1.30, "AFTER", TEAL)
    add_text(slide, 7.82, 2.68, 4.48, 0.52, "X′ᵣ ∈ ℝᴮ×ˢ×ᴴ⁄ᴾ×ᴰ", 20, INK, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 8.02, 3.43, 4.0, 1.35, "rank r 拥有：\n• 完整 sequence\n• 一组 attention heads", 17, MUTED)
    add_text(slide, 7.82, 5.05, 4.48, 0.34, "[1, 37,760, 14, 128]", 17, TEAL, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 0.82, 6.18, 11.72, 0.32, "元素守恒：(S/P)·H·D = S·(H/P)·D；Ulysses 不是复制，而是 ownership 重映射。", 15, INK, True, FONT, PP_ALIGN.CENTER)
    add_footer(slide, 12, "attention/parallel/ulysses.py:320–334；distributed/comm.py:36–62")

    # 13 — why three
    slide = blank_slide(prs)
    add_title(slide, "12 · Three collectives", "为什么 Q、K、V 要做三次 All-to-All", "三个 tensor 都必须从 sequence-owner 布局变成 head-owner 布局")
    rows = [
        ("Q", "获得其他 rank 的 query tokens", "每个 head-owner rank 要输出所有 token 的本地 head 子集", PURPLE),
        ("K", "获得完整 key sequence", "每个 query 必须与全局 keys 交互", BLUE),
        ("V", "获得完整 value sequence", "softmax 权重需要聚合全局 values", ORANGE),
    ]
    for i, (label, reason, consequence, color) in enumerate(rows):
        y = 1.78 + i * 1.25
        add_pill(slide, 0.78, y + 0.24, 0.72, label, color, BG, 15)
        add_rect(slide, 1.72, y, 10.70, 0.96, PANEL, GRID)
        add_text(slide, 2.02, y + 0.15, 3.66, 0.30, reason, 15, color, True)
        add_text(slide, 5.78, y + 0.14, 6.10, 0.48, consequence, 14, INK)
    add_rect(slide, 0.78, 5.76, 11.64, 0.74, PANEL_2, GRID)
    add_text(slide, 1.02, 5.94, 2.58, 0.28, "可否融合为一次？", 14, TEAL, True)
    add_text(slide, 3.62, 5.88, 8.46, 0.40, "H3 形状上可评估 pack(Q,K,V)，但总字节数不变；收益主要来自减少 launch/sync，且需承担 pack/unpack copy。", 12, MUTED)
    add_footer(slide, 13, "attention/parallel/ulysses.py:331–334")

    # 14 — head ownership
    slide = blank_slide(prs)
    add_title(slide, "13 · Local attention", "A2A 后：每张卡计算完整 S 的 14 个 heads", "Flash/CuDNN attention 不需要物化 S×S score matrix")
    colors = [TEAL, BLUE, ORANGE, PURPLE]
    for i in range(4):
        x = 0.72 + i * 3.02
        add_rect(slide, x, 1.88, 2.72, 3.60, PANEL, colors[i])
        add_pill(slide, x + 0.18, 2.10, 0.72, f"R{i}", colors[i], BG, 11)
        lo, hi = i * 14, (i + 1) * 14 - 1
        add_text(slide, x + 0.20, 2.72, 2.30, 0.40, f"heads {lo}–{hi}", 17, INK, True, MONO, PP_ALIGN.CENTER)
        add_text(slide, x + 0.20, 3.42, 2.30, 0.58, "完整 packed S\n= 37,760", 16, colors[i], True, MONO, PP_ALIGN.CENTER)
        add_text(slide, x + 0.20, 4.42, 2.30, 0.46, "QKᵀ → softmax → V", 13, MUTED, font=MONO, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.72, 5.86, 11.78, 0.58, PANEL_2, GRID)
    add_text(slide, 0.98, 6.02, 11.25, 0.25, "全局计算量基本不变；Ulysses 将 head 维并行化，并把 activation/attention 工作均匀分布到 4 张卡。", 14, INK, True, FONT, PP_ALIGN.CENTER)
    add_footer(slide, 14, "attention/layer.py:280–288；H3 num_attention_heads=56")

    # 15 — reverse
    slide = blank_slide(prs)
    add_title(slide, "14 · Reverse reshard", "为什么 Attention 输出还要再做一次 All-to-All", "下一步 OutProj、残差与 MLP 需要恢复 sequence-sharded、完整 hidden 的布局")
    add_process_box(slide, 0.72, 2.00, 3.22, 2.72, "Attention output", "[1, S, 14, 128]\n\n完整 sequence\n局部 heads", TEAL, 18, 16)
    add_chevron(slide, 4.20, 2.96, 0.70, 0.76, ORANGE)
    add_process_box(slide, 5.14, 2.00, 3.22, 2.72, "Reverse A2A", "scatter sequence\ngather heads\n\n第 4 次 collective", ORANGE, 18, 15)
    add_chevron(slide, 8.61, 2.96, 0.70, 0.76, PURPLE)
    add_process_box(slide, 9.56, 2.00, 3.10, 2.72, "Local continuation", "[S/4, 56, 128]\n→ [S/4, 5,376]\n\nOutProj + MLP", PURPLE, 18, 14)
    add_rect(slide, 0.72, 5.24, 11.94, 0.94, PANEL, GRID)
    add_text(slide, 0.98, 5.49, 11.42, 0.42, "如果不恢复布局，下一层的 sequence-local MLP 会变成 head-sharded hidden，残差连接和 RowParallelLinear 语义都不再成立。", 14, INK, True, FONT, PP_ALIGN.CENTER)
    add_footer(slide, 15, "attention/parallel/ulysses.py:414–473；minimax_h3_transformer.py:457–458")

    # 16 — counts
    slide = blank_slide(prs)
    add_title(slide, "15 · Communication count", "从一个 block 放大到一次请求", "计数来自源码循环；不包含 warmup、Text Encoder 与 VAE")
    add_metric(slide, 0.62, 1.78, 2.86, "4", "All-to-All / block", ORANGE, "Q + K + V + Output")
    add_chevron(slide, 3.58, 2.10, 0.38, 0.48, TEAL)
    add_metric(slide, 4.08, 1.78, 2.86, "200", "All-to-All / forward", TEAL, "4 × 50 blocks")
    add_chevron(slide, 7.05, 2.10, 0.38, 0.48, PURPLE)
    add_metric(slide, 7.56, 1.78, 2.86, "800", "All-to-All / request", PURPLE, "200 × 4 transitions")
    add_metric(slide, 10.62, 1.78, 2.05, "4", "final AllGather", RED, "1 × transition")
    add_rect(slide, 0.62, 3.38, 12.05, 2.56, PANEL, GRID)
    add_text(slide, 0.94, 3.72, 5.08, 0.38, "采样步数的精确定义", 17, TEAL, True)
    add_text(slide, 0.94, 4.30, 5.20, 0.92, "num_inference_steps = 5\n生成 5 个 sigma points\n循环 len(sigmas)-1 = 4 次", 16, INK, font=MONO)
    add_text(slide, 6.62, 3.72, 5.10, 0.38, "为何不是 1,000 次 A2A？", 17, ORANGE, True)
    add_text(slide, 6.62, 4.30, 5.20, 0.92, "H3 是 CFG-distilled：每个 transition 只执行一次 transformer forward，不做正/负条件双分支。", 15, MUTED)
    add_footer(slide, 16, "denoise_loop.py；MiniMaxH3DiTModel._cache_dit_adapter_config.has_separate_cfg=False")

    # 17 — volume model
    slide = blank_slide(prs)
    add_title(slide, "16 · Analytical traffic model", "U4 下 Attention A2A 的理论通信体量", "BF16；按每 rank 远端发送量估算，receive volume 相同")
    add_rect(slide, 0.62, 1.70, 12.05, 1.02, PANEL, GRID)
    add_text(slide, 0.88, 1.91, 11.55, 0.26, "local tensor = (S/P)·H·D·2 B = 9,440·56·128·2 = 129.06 MiB", 16, INK, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 0.88, 2.27, 11.55, 0.24, "remote fraction ≈ (P−1)/P = 3/4  →  96.80 MiB / collective / rank", 14, TEAL, True, MONO, PP_ALIGN.CENTER)
    add_metric(slide, 0.62, 3.16, 2.86, "96.8 MiB", "one A2A", BLUE, "remote send / rank")
    add_metric(slide, 3.65, 3.16, 2.86, "387.2 MiB", "one block", ORANGE, "4 collectives")
    add_metric(slide, 6.68, 3.16, 2.86, "18.9 GiB", "one forward", TEAL, "50 blocks")
    add_metric(slide, 9.71, 3.16, 2.96, "75.6 GiB", "one request", PURPLE, "4 transitions")
    add_rect(slide, 0.62, 4.68, 12.05, 1.23, PANEL_2, GRID)
    add_text(slide, 0.90, 4.94, 2.52, 0.30, "解释边界", 14, RED, True)
    add_text(slide, 3.28, 4.87, 8.98, 0.58, "这是尺寸模型而非链路计数器：实际 NCCL 协议、NVLink/NVSwitch 拓扑、算法与 overlap 会改变时延，但不会改变必须重分布的逻辑 payload。", 13, MUTED)
    add_footer(slide, 17, "Derived from S=37,760, P=4, H=56, D=128, BF16")

    # 18 — final allgather
    slide = blank_slide(prs)
    add_title(slide, "17 · Exit boundary", "为什么 50 层之后还有一个大 AllGather", "当前架构：DiT 内部 sharded，denoise state 在 DiT 外部 replicated")
    add_process_box(slide, 0.62, 1.88, 3.36, 2.56, "Block stack output", "每 rank\n[S/4, 5,376]\n= [9,440, 5,376]", TEAL, 17, 15)
    add_chevron(slide, 4.23, 2.72, 0.62, 0.72, RED)
    add_process_box(slide, 5.08, 1.88, 3.36, 2.56, "sp_gather", "AllGather(dim=0)\n\n每 rank 得到\n[S, 5,376]", RED, 17, 15)
    add_chevron(slide, 8.69, 2.72, 0.62, 0.72, ORANGE)
    add_process_box(slide, 9.54, 1.88, 3.13, 2.56, "Replicated tail", "final layer\nunpack logits\nscheduler update\nnext-step state", ORANGE, 17, 13)
    add_metric(slide, 0.62, 4.86, 3.50, "387.2 MiB", "materialized hidden / rank", RED, "S·5376·2 bytes")
    add_metric(slide, 4.39, 4.86, 3.50, "≈290.4 MiB", "ring recv / rank", ORANGE, "(P−1)/P estimate")
    add_rect(slide, 8.16, 4.86, 4.51, 1.16, PANEL, GRID)
    add_text(slide, 8.42, 5.07, 3.95, 0.28, "为什么是 AllGather 而不是 Gather？", 14, TEAL, True)
    add_text(slide, 8.42, 5.45, 3.95, 0.30, "每个 rank 都继续维护完整 denoise state。", 13, INK, True)
    add_footer(slide, 18, "minimax_h3_transformer.py:1123–1129；hooks/sequence_parallel.py:518–532")

    # 19 — optimization
    slide = blank_slide(prs)
    add_title(slide, "18 · Strict boundary optimization", "先本地 final layer，再 AllGather 紧凑 logits", "通信宽度 5,376 → 96+32=128，逻辑 payload 缩小 42×")
    add_rect(slide, 0.62, 1.72, 5.74, 3.82, PANEL, RED)
    add_pill(slide, 0.92, 2.02, 1.54, "CURRENT", RED, BG)
    add_text(slide, 0.94, 2.68, 5.10, 1.52, "[S/4, 5,376]\n↓ AllGather\n[S, 5,376]\n↓ final layer\n[S, 128]", 19, INK, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 0.94, 4.67, 5.10, 0.36, "materialized = 387.2 MiB / rank", 15, RED, True, MONO, PP_ALIGN.CENTER)
    add_chevron(slide, 6.51, 3.15, 0.42, 0.62, TEAL)
    add_rect(slide, 7.15, 1.72, 5.52, 3.82, PANEL_2, TEAL)
    add_pill(slide, 7.45, 2.02, 1.78, "OPTIMIZED", TEAL, BG)
    add_text(slide, 7.47, 2.68, 4.90, 1.52, "[S/4, 5,376]\n↓ local final layer\n[S/4, 128]\n↓ AllGather\n[S, 128]", 19, INK, True, MONO, PP_ALIGN.CENTER)
    add_text(slide, 7.47, 4.67, 4.90, 0.36, "materialized = 9.22 MiB / rank", 15, TEAL, True, MONO, PP_ALIGN.CENTER)
    add_rect(slide, 0.62, 5.85, 12.05, 0.58, PANEL, GRID)
    add_text(slide, 0.88, 6.00, 11.55, 0.25, "该优化已存在于后续提交 7b76b6446（#6173），但不在当前 f76f8e58fb 分支；Attention 内 800 次 A2A 不受影响。", 12, MUTED, font=MONO, align=PP_ALIGN.CENTER)
    add_footer(slide, 19, "Later commit: 7b76b6446 · Optimize MiniMax-H3 strict Ulysses boundaries (#6173)")

    # 20 — optimization priorities
    slide = blank_slide(prs)
    add_title(slide, "19 · Optimization agenda", "后续优化应按 payload、频次与实现风险排序", "必须以固定输入与正确性校验验证，不接受质量回归")
    headers = ["候选项", "主要收益", "通信字节", "实现风险", "优先级"]
    widths = [3.10, 3.05, 1.70, 1.65, 1.45]
    x0, y0 = 0.68, 1.72
    x = x0
    for h, w in zip(headers, widths):
        add_rect(slide, x, y0, w, 0.58, TEAL, BG, radius=False)
        add_text(slide, x + 0.06, y0 + 0.13, w - 0.12, 0.24, h, 12, BG, True, align=PP_ALIGN.CENTER)
        x += w
    rows = [
        ("严格 SP 出口边界", "final layer 本地化", "42×↓", "低—中", "P0", TEAL),
        ("融合 QKV A2A", "减少 launch/sync", "基本不变", "中", "P1", ORANGE),
        ("通信—计算 overlap", "隐藏 exposed latency", "不变", "中—高", "P1", BLUE),
        ("端到端 sharded denoise", "移除末端 AllGather", "进一步↓", "高", "P2", PURPLE),
        ("局部 embedding/RoPE", "减少 replicated work", "间接", "中", "P1", GREEN),
    ]
    for i, row in enumerate(rows):
        y = y0 + 0.58 + i * 0.77
        x = x0
        for j, (value, w) in enumerate(zip(row[:5], widths)):
            fill = PANEL if i % 2 == 0 else PANEL_2
            add_rect(slide, x, y, w, 0.77, fill, GRID, radius=False)
            color = row[5] if j in (0, 4) else INK
            add_text(slide, x + 0.08, y + 0.17, w - 0.16, 0.35, value, 11, color, j in (0, 4), align=PP_ALIGN.CENTER)
            x += w
    add_text(slide, 0.72, 6.45, 11.8, 0.26, "P0 优先移除明显过宽的边界通信；QKV fusion 只有在 profiler 证明 launch-bound 时才值得推进。", 13, MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 20, "优化原则：先消除不必要 work，再优化必要 collective")

    # 21 — profiling
    slide = blank_slide(prs)
    add_title(slide, "20 · Profiling protocol", "如何用 Nsight Systems 验证 Ulysses 热路径", "Linux/B300 采集 .nsys-rep；Mac 使用 Nsight Systems GUI 打开分析")
    add_process_box(slide, 0.62, 1.74, 3.42, 3.98, "1 · Fixed workload", "固定：prompt / seed / H/W/T / steps / backend / U4\n\n性能 profile 建议 2 steps；端到端报告使用真实 steps。", TEAL, 17, 14)
    add_process_box(slide, 4.29, 1.74, 4.02, 3.98, "2 · Nsight Systems", "观察：\n• NCCL AllToAll 时间线\n• kernel 与 collective gap\n• exposed communication\n• block 周期是否稳定\n• final AllGather payload", ORANGE, 17, 14)
    add_process_box(slide, 8.56, 1.74, 4.11, 3.98, "3 · Source attribution", "期望计数：\n• 4 A2A / block\n• 200 A2A / forward\n• 800 A2A / request\n• 4 final AllGather\n\n用 NVTX/call stack 映射源码。", PURPLE, 17, 13)
    add_rect(slide, 0.62, 6.02, 12.05, 0.48, PANEL_2, GRID)
    add_text(slide, 0.88, 6.13, 11.55, 0.24, "nsys 回答系统时间线；ncu 仅用于单个 attention/GEMM kernel 的深层硬件指标，两者不能互相替代。", 13, INK, True, FONT, PP_ALIGN.CENTER)
    add_footer(slide, 21, "profiling/minimax-h3-sp8-nsys-20260820/trace/nsys_smoke.nsys-rep（现有 trace 为 U8，仅供同构路径参考）")

    # 22 — validation
    slide = blank_slide(prs)
    add_title(slide, "21 · Evaluation design", "优化前后如何做学术上可信的对比", "分离 correctness、wall clock、communication 与 kernel 指标")
    columns = [
        ("Correctness", ["相同 seed/input", "输出 shape", "音视频 hash 或容差", "无 NaN/同步错误"], TEAL),
        ("Performance", ["DiT / diffuse latency", "E2E latency", "GPU peak memory", "至少复现一次"], ORANGE),
        ("Communication", ["collective 次数", "payload bytes", "exposed NCCL time", "overlap ratio"], PURPLE),
        ("Attribution", ["相同 backend", "相同 U4 topology", "相同 compile state", "source→trace 对齐"], BLUE),
    ]
    for i, (title, bullets, color) in enumerate(columns):
        x = 0.62 + i * 3.03
        add_rect(slide, x, 1.78, 2.86, 4.72, PANEL if i % 2 == 0 else PANEL_2, color)
        add_pill(slide, x + 0.20, 2.03, 2.10, title.upper(), color, BG, 10)
        for j, item in enumerate(bullets):
            add_rect(slide, x + 0.25, 2.78 + j * 0.74, 0.07, 0.07, color, color, radius=False)
            add_text(slide, x + 0.43, 2.64 + j * 0.74, 2.10, 0.42, item, 13, INK)
        add_text(slide, x + 0.24, 5.92, 2.32, 0.28, f"CHECK {i+1}/4", 9, color, True, MONO, PP_ALIGN.RIGHT)
    add_footer(slide, 22, "Benchmark discipline: same model/input/shape/steps/parallel config")

    # 23 — conclusions and references
    slide = blank_slide(prs)
    add_title(slide, "22 · Takeaways", "结论与源码索引", "Ulysses 的本质是 ownership 转换；性能分析必须同时看边界和层内通信")
    add_rect(slide, 0.62, 1.72, 6.02, 4.86, PANEL, GRID)
    add_text(slide, 0.92, 2.00, 5.35, 0.34, "最终结论", 18, TEAL, True)
    add_bullets(slide, 0.95, 2.55, 5.18, [
        "H3 用 _sp_plan 将 50 层 DiT 包在一个 SP dynamic scope 内。",
        "每层通过 3×QKV A2A 获得 full-S/head-shard，再通过 1×Output A2A 恢复 S-shard。",
        "U4 请求的主要重复通信是 800 次 A2A；末端另有 4 次 AllGather。",
        "先本地 final layer 可把末端 gather 宽度从 5,376 降到 128；进一步消除 gather 需要端到端 sharded denoise。",
    ], 14, gap=0.82)
    add_rect(slide, 6.90, 1.72, 5.77, 4.86, PANEL_2, GRID)
    add_text(slide, 7.20, 2.00, 5.12, 0.34, "关键源码索引", 18, ORANGE, True)
    refs = [
        "minimax_h3_transformer.py:785–804, 1109–1129",
        "minimax_h3_transformer.py:328–459, 617–681",
        "attention/layer.py:263–290",
        "attention/parallel/ulysses.py:319–334, 414–473",
        "distributed/comm.py:16–98",
        "hooks/sequence_parallel.py:236–268, 518–532",
        "results/h3-b300-tp1-u4-steps5/summary.json",
    ]
    for i, ref in enumerate(refs):
        add_text(slide, 7.22, 2.58 + i * 0.48, 5.00, 0.30, f"{i+1:02d}  {ref}", 10, INK, font=MONO)
    add_pill(slide, 8.56, 6.00, 2.40, "END · DISCUSSION", TEAL)
    add_footer(slide, 23, "Local evidence package · branch f76f8e58fb · later optimization 7b76b6446")

    prs.save(OUT_FILE)
    return OUT_FILE, len(prs.slides)


if __name__ == "__main__":
    output, count = build_deck()
    print(f"wrote {output}")
    print(f"slides {count}")
